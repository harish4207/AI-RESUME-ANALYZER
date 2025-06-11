import os
from pptx import Presentation
from pptx.util import Inches

def fill_aicte_template(template_path, output_path):
    """Fill the AICTE project template with our AI Resume Analyzer project details."""
    try:
        prs = Presentation(template_path)
        
        # Project details
        project_data = {
            "problem_statement": """
            In today's highly competitive job market, a significant challenge for job seekers is the widespread adoption of Applicant Tracking Systems (ATS) by companies for initial resume screening. 
            These automated systems often filter out qualified candidates whose resumes are not optimized for ATS algorithms, leading to a substantial loss of talent. 
            The current manual process of resume creation and optimization is time-consuming and often ineffective, resulting in many deserving individuals failing to reach the human recruiter stage. 
            There is a critical need for an intelligent solution that bridges this gap, enabling job seekers to create ATS-friendly resumes that effectively highlight their skills and experience, thereby increasing their chances of securing interviews.
            """,
            "proposed_system_solution": """
            The AI Resume Analyzer is an AI-powered web application designed to optimize resumes for Applicant Tracking Systems (ATS) and specific job requirements. This system provides a comprehensive solution to enhance resume visibility and improve interview chances through the following components:

            - **Data Collection & Extraction**:
                - Gather resume data (PDF, DOCX, TXT) and job descriptions.
                - Utilize PyMuPDF and python-docx for text extraction.

            - **Data Preprocessing**:
                - Clean and preprocess collected text (lowercasing, punctuation, stop words, lemmatization).

            - **Core Algorithms**:
                - Employ spaCy for keyword extraction (skills, responsibilities, industry terms).
                - Implement rule-based ATS scoring (format, contact, keyword density, readability).
                - Apply TF-IDF and Cosine Similarity for resume-job description matching, highlighting missing keywords.

            - **Deployment**:
                - Develop a user-friendly Streamlit web application.
                - Deploy on a scalable platform (e.g., AWS, Azure, GCP).

            - **Evaluation & Feedback**:
                - Assess performance via ATS scores and job match percentages.
                - Provide detailed, categorized suggestions for resume improvement.
            """,
            "system_development_approach": """
            The AI Resume Analyzer is developed with a focus on modularity, scalability, and user-friendliness. The overall strategy involves leveraging robust open-source libraries and a streamlined development methodology.

            - **System Requirements**:
                - **Hardware**: A standard computer with at least 8GB RAM (16GB recommended) and multi-core processor for optimal performance.
                - **Operating System**: Cross-platform compatibility (Windows, macOS, Linux).
                - **Software**: Python 3.8+ environment, Streamlit framework, and necessary Python packages.

            - **Libraries Required to Build the Model**:
                - `streamlit`: For building the interactive web application interface.
                - `spacy`: For Natural Language Processing (NLP) tasks, including text parsing, tokenization, and keyword extraction. Requires `en_core_web_sm` model.
                - `scikit-learn`: For machine learning functionalities, specifically `TfidfVectorizer` and `cosine_similarity` for job description matching.
                - `python-docx`: For extracting text content from .docx (Microsoft Word) resume files.
                - `PyMuPDF` (fitz): For extracting text content from .pdf (Portable Document Format) resume files.
                - `pandas`: For data manipulation and structuring, especially for charting.
                - `plotly.graph_objects` & `plotly.express`: For creating interactive data visualizations (gauge charts, bar charts) within the Streamlit application.
                - `re`: Python's built-in regular expression module for advanced text pattern matching and validation.
            """,
            "algorithm_deployment": """
            **Algorithms Used:**
            - **Text Preprocessing**: Custom cleaning functions (e.g., lowercasing, punctuation removal, stop word removal) for raw resume and job description text.
            - **Keyword Extraction**: spaCy's NLP capabilities are used to identify relevant nouns, verbs, and entities, forming a robust set of keywords.
            - **ATS Scoring Logic**: Rule-based algorithms evaluate resume sections (format, keywords, contact info, readability) against predefined criteria.
            - **Similarity Calculation**: TF-IDF (Term Frequency-Inverse Document Frequency) vectorization combined with Cosine Similarity is used to quantify the match between a resume and a job description.
            
            **Deployment:**
            The application is deployed as a Streamlit web application, allowing easy access via a web browser. 
            It is designed to be scalable and can be hosted on cloud platforms (e.g., AWS, Azure, GCP) or on-premise servers, ensuring broad accessibility for users.
            """,
            "result_output_image": "[INSERT SCREENSHOT OF THE APP HERE]", # Placeholder for an image
            "conclusion": """
            The AI Resume Analyzer successfully provides job seekers with a powerful tool to optimize their resumes for ATS and improve their chances of securing interviews. 
            By offering detailed, data-driven suggestions, the system empowers users to create highly effective resumes that align with industry standards and specific job requirements. 
            The project demonstrates the practical application of AI and NLP in addressing real-world career challenges.
            """,
            "future_scope": """
            Future enhancements for the AI Resume Analyzer include:
            - **Cover Letter Analysis**: Integrate functionality to analyze and provide suggestions for cover letters.
            - **Interview Preparation**: Add features for common interview questions and mock interview simulations.
            - **Personalized Learning Paths**: Suggest courses or certifications based on skill gaps identified.
            - **Multi-language Support**: Expand support for resume analysis in various languages.
            - **User Authentication & Profile Management**: Allow users to save their resumes and track progress over time.
            - **AI-powered Resume Generation**: Offer a feature to generate resumes from scratch based on user input and job descriptions.
            """,
            "references": [
                "AICTE Guidelines for Project Development",
                "Industry Best Practices for Resume Optimization",
                "Research Papers on ATS Systems and Resume Analysis",
                "Technical Documentation for Streamlit, spaCy, scikit-learn, PyMuPDF, python-docx"
            ]
        }
        
        # Fill each slide with project data based on placeholder text or slide titles
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_lower = run.text.lower()
                            
                            if "problem statement" in text_lower and "solution" not in text_lower:
                                run.text = project_data["problem_statement"]
                            elif "proposed system/solution" in text_lower:
                                run.text = project_data["proposed_system_solution"]
                            elif "system development approach" in text_lower or "technology used" in text_lower:
                                run.text = project_data["system_development_approach"]
                            elif "algorithm & deployment" in text_lower:
                                run.text = project_data["algorithm_deployment"]
                            elif "result (output image)" in text_lower:
                                run.text = project_data["result_output_image"]
                                # You'll need to manually add an image here using pptx API if it's a picture placeholder
                                # For now, it just replaces the text.
                            elif "conclusion" in text_lower:
                                run.text = project_data["conclusion"]
                            elif "future scope" in text_lower:
                                run.text = project_data["future_scope"]
                            elif "references" in text_lower:
                                run.text = "\n".join(f"• {ref}" for ref in project_data["references"])

        # Save the filled template
        prs.save(output_path)
        print(f"Template filled successfully! Saved to: {output_path}")
        
    except Exception as e:
        print(f"Error filling template: {str(e)}")

def main():
    # Template and output paths
    template_path = r"C:\Users\HARISH\OneDrive\المستندات\Project_template_AICTE.pptx"
    output_path = "AI_Resume_Analyzer_AICTE_Project_Filled.pptx"
    
    # Check if template exists
    if not os.path.exists(template_path):
        print(f"Error: Template file not found at {template_path}")
        return
    
    # Fill the template
    print("Filling AICTE project template...")
    fill_aicte_template(template_path, output_path)

if __name__ == "__main__":
    main()
